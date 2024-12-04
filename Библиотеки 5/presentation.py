slides ={'seed':'Initialize the random number generator.', 
         'getstate':'Return an object capturing the current internal state of the generator. This object can be passed to setstate() to restore the state.', 
         'getrandbits':'Returns a Python integer with k random bits. This method is supplied with the MersenneTwister generator and some other generators may also provide it as an optional part of the API. When available, getrandbits() enables randrange() to handle arbitrarily large ranges.',
         'randrange':'Return a randomly selected element from range(start, stop, step). This is equivalent to choice(range(start, stop, step)), but doesn’t actually build a range object.',
         'randint':'Return a random integer N such that a <= N <= b. Alias for randrange(a, b+1).',
         'choice':'Return a random element from the non-empty sequence seq. If seq is empty, raises IndexError.',
         'choices':'Return a k sized list of elements chosen from the population with replacement. If the population is empty, raises IndexError.',
         'shuffle':'Shuffle the sequence x in place.',
         'sample':'Return a k length list of unique elements chosen from the population sequence or set. Used for random sampling without replacement.',
         'random':'Return the next random floating point number in the range [0.0, 1.0).',
         'uniform':'Return a random floating point number N such that a <= N <= b for a <= b and b <= N <= a for b < a.',
         'triangular(low, high, mode)':'Return a random floating point number N such that low <= N <= high and with the specified mode between those bounds. The low and high bounds default to zero and one. The mode argument defaults to the midpoint between the bounds, giving a symmetric distribution.',
         'betavariate':'Beta distribution. Conditions on the parameters are alpha > 0 and beta > 0. Returned values range between 0 and 1.',
         'expovariate':'Exponential distribution. lambd is 1.0 divided by the desired mean. It should be nonzero. (The parameter would be called “lambda”, but that is a reserved word in Python.) Returned values range from 0 to positive infinity if lambd is positive, and from negative infinity to 0 if lambd is negative.',
         'gammavariate':'Gamma distribution. (Not the gamma function!) Conditions on the parameters are alpha > 0 and beta > 0.',
         'gauss':'Gaussian distribution. mu is the mean, and sigma is the standard deviation. This is slightly faster than the normalvariate() function defined below.',
         'lognormvariate':'Log normal distribution. If you take the natural logarithm of this distribution, you’ll get a normal distribution with mean mu and standard deviation sigma. mu can have any value, and sigma must be greater than zero.',
         'normalvariate':'Normal distribution. mu is the mean, and sigma is the standard deviation.',
         'vonmisesvariate':'mu is the mean angle, expressed in radians between 0 and 2*pi, and kappa is the concentration parameter, which must be greater than or equal to zero. If kappa is equal to zero, this distribution reduces to a uniform random angle over the range 0 to 2*pi.',
         'paretovariate':'Pareto distribution. alpha is the shape parameter.',
         'weibullvariate':'Weibull distribution. alpha is the scale parameter and beta is the shape parameter.',
         'setstate':'state should have been obtained from a previous call to getstate(), and setstate() restores the internal state of the generator to what it was at the time getstate() was called.',
         'triangular':'Return a random floating point number N such that low <= N <= high and with the specified mode between those bounds. The low and high bounds default to zero and one. The mode argument defaults to the midpoint between the bounds, giving a symmetric distribution.'}

import pptx
import random
pres = pptx.Presentation()
methods = dir(random)[39:]
lay = pres.slide_layouts[0]
for method in methods:
    slide = pres.slides.add_slide(lay)
    title = slide.shapes.title
    box = slide.shapes.placeholders[0].text_frame
    title.text = method
    text = box.add_paragraph()
    text.text = slides.get(method, '')
    text.font.name = "Courier New"
    text.font.size = pptx.util.Pt(14)
pres.save("res.pptx")
